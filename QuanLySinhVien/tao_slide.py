# -*- coding: utf-8 -*-
"""
Script tạo file PPTX thuyết trình về các thuật toán trong hệ thống
Quản lý Sinh viên - Môn Kỹ thuật Lập trình Nâng cao
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# CẤU HÌNH MÀU SẮC
# ============================================================
COLOR_PRIMARY = RGBColor(0x1A, 0x23, 0x7E)       # Xanh đậm
COLOR_SECONDARY = RGBColor(0x30, 0x49, 0xC7)     # Xanh vừa
COLOR_ACCENT = RGBColor(0xFF, 0x6F, 0x00)        # Cam
COLOR_SUCCESS = RGBColor(0x2E, 0x7D, 0x32)       # Xanh lá
COLOR_DANGER = RGBColor(0xC6, 0x28, 0x28)        # Đỏ
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK = RGBColor(0x21, 0x21, 0x21)
COLOR_GRAY = RGBColor(0x61, 0x61, 0x61)
COLOR_LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
COLOR_CODE_BG = RGBColor(0x26, 0x32, 0x38)
COLOR_PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
COLOR_TEAL = RGBColor(0x00, 0x69, 0x5C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================================
# HÀM TIỆN ÍCH
# ============================================================

def add_background(slide, color=COLOR_WHITE):
    """Thêm nền cho slide"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, radius=None):
    """Thêm hình chữ nhật bo góc làm nền"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_rect(slide, left, top, width, height, color):
    """Thêm hình chữ nhật"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=COLOR_DARK, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Segoe UI'):
    """Thêm hộp văn bản"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(tf, items, font_size=16, color=COLOR_DARK, 
                              spacing=Pt(6), font_name='Segoe UI', bold_first=False):
    """Thêm nội dung dạng bullet vào text frame"""
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        if bold_first and ':' in item:
            # Cannot do partial bold easily, just bold entire line
            pass


def add_header_bar(slide, title_text):
    """Thêm thanh tiêu đề phía trên"""
    # Thanh gradient trên cùng
    bar = add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), COLOR_PRIMARY)
    # Tiêu đề
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.8),
                 title_text, font_size=30, color=COLOR_WHITE, bold=True)
    # Đường accent dưới header
    add_rect(slide, Inches(0), Inches(1.1), prs.slide_width, Inches(0.06), COLOR_ACCENT)


def add_section_title(slide, left, top, text, color=COLOR_SECONDARY, font_size=22):
    """Thêm tiêu đề mục"""
    add_text_box(slide, left, top, Inches(11), Inches(0.5), text,
                 font_size=font_size, color=color, bold=True)
    # Đường gạch chân
    add_rect(slide, left, top + Inches(0.45), Inches(2.5), Inches(0.04), color)


def add_code_block(slide, left, top, width, height, code_text, font_size=12):
    """Thêm khung code"""
    bg = add_shape_bg(slide, left, top, width, height, COLOR_CODE_BG)
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                      width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = code_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0xA5, 0xD6, 0xA7)
        p.font.name = 'Consolas'
        p.space_after = Pt(2)


def add_complexity_table(slide, left, top, data, col_widths=None):
    """Thêm bảng độ phức tạp"""
    rows = len(data)
    cols = len(data[0])
    if col_widths is None:
        col_widths = [Inches(2.5)] * cols
    
    table_shape = slide.shapes.add_table(rows, cols, left, top,
                                          sum(col_widths), Inches(0.4 * rows))
    table = table_shape.table
    
    for c, w in enumerate(col_widths):
        table.columns[c].width = w
    
    for r, row_data in enumerate(data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = 'Segoe UI'
                paragraph.alignment = PP_ALIGN.CENTER
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = COLOR_WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_PRIMARY
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xEA, 0xF6)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_WHITE
    
    return table_shape


def add_info_card(slide, left, top, width, height, title, content, title_color=COLOR_SECONDARY):
    """Thêm thẻ thông tin"""
    card = add_shape_bg(slide, left, top, width, height, RGBColor(0xF8, 0xF9, 0xFA))
    # Border top
    add_rect(slide, left, top, width, Inches(0.05), title_color)
    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.4), Inches(0.4),
                 title, font_size=15, color=title_color, bold=True)
    # Content
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.5),
                                      width - Inches(0.4), height - Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


# ============================================================
# SLIDE 1: TRANG BÌA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_background(slide, COLOR_PRIMARY)

# Decorative shapes
add_rect(slide, Inches(0), Inches(0), Inches(0.15), prs.slide_height, COLOR_ACCENT)
add_shape_bg(slide, Inches(8.5), Inches(-1), Inches(6), Inches(3), 
             RGBColor(0x28, 0x3C, 0x8E))
add_shape_bg(slide, Inches(9), Inches(5.5), Inches(5.5), Inches(3),
             RGBColor(0x28, 0x3C, 0x8E))

# Tên trường
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
             "TRƯỜNG ĐẠI HỌC HÙNG VƯƠNG", font_size=18, color=RGBColor(0xBB, 0xDE, 0xFB), bold=False)

# Tên đề tài
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.6),
             "TIỂU LUẬN MÔN HỌC", font_size=20, color=COLOR_ACCENT, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.6),
             "KỸ THUẬT LẬP TRÌNH NÂNG CAO", font_size=36, color=COLOR_WHITE, bold=True)

# Đường kẻ
add_rect(slide, Inches(0.8), Inches(2.55), Inches(4), Inches(0.06), COLOR_ACCENT)

# Tên đề tài chính
add_text_box(slide, Inches(0.8), Inches(2.9), Inches(11), Inches(1.2),
             "Xây dựng phần mềm Quản lý Sinh viên\nsử dụng cấu trúc Mảng động và Giao diện đồ họa",
             font_size=26, color=RGBColor(0xE3, 0xF2, 0xFD), bold=False)

# Thông tin GVHD & SVTH
info_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(6), Inches(2.2))
tf = info_box.text_frame
tf.word_wrap = True

lines = [
    ("GVHD:  ", "TS. Nghiêm Văn Tính"),
    ("SVTH:  ", "Đặng Đình Đạt"),
]

for label, value in lines:
    p = tf.add_paragraph()
    run1 = p.add_run()
    run1.text = label
    run1.font.size = Pt(20)
    run1.font.color.rgb = COLOR_ACCENT
    run1.font.bold = True
    run1.font.name = 'Segoe UI'
    
    run2 = p.add_run()
    run2.text = value
    run2.font.size = Pt(20)
    run2.font.color.rgb = COLOR_WHITE
    run2.font.bold = False
    run2.font.name = 'Segoe UI'
    p.space_after = Pt(8)

# Năm
add_text_box(slide, Inches(0.8), Inches(6.7), Inches(4), Inches(0.5),
             "Năm học 2025 - 2026", font_size=16, color=RGBColor(0x90, 0xCA, 0xF9))


# ============================================================
# SLIDE 2: MỤC LỤC
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "📋  MỤC LỤC THUYẾT TRÌNH")

items = [
    ("01", "Tổng quan đề tài", "Giới thiệu bài toán và phạm vi"),
    ("02", "Kiến trúc hệ thống", "Mô hình MVC, cấu trúc module"),
    ("03", "Cấu trúc dữ liệu Mảng động", "Nguyên lý, Doubling Strategy, Amortized Analysis"),
    ("04", "Thuật toán Thêm / Xóa phần tử", "Thêm cuối O(1)*, Xóa O(n) với dịch chuyển"),
    ("05", "Thuật toán Tìm kiếm tuyến tính", "Linear Search trên nhiều trường, O(n)"),
    ("06", "Thuật toán Sắp xếp (TimSort)", "Hybrid Merge Sort + Insertion Sort, O(n log n)"),
    ("07", "Thuật toán Lọc & Thống kê", "Filter, Aggregate, Set-based deduplication"),
    ("08", "Thuật toán Phân trang", "Slice-based pagination, O(k)"),
    ("09", "Bảng tổng hợp độ phức tạp", "So sánh tất cả thuật toán"),
    ("10", "Kết luận", "Đánh giá và hướng phát triển"),
]

for i, (num, title, desc) in enumerate(items):
    y = Inches(1.5) + Inches(i * 0.55)
    # Số thứ tự
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLOR_PRIMARY if i % 2 == 0 else COLOR_SECONDARY
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = num
    circle.text_frame.paragraphs[0].font.size = Pt(14)
    circle.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    circle.text_frame.paragraphs[0].font.name = 'Segoe UI'
    
    # Tiêu đề
    add_text_box(slide, Inches(1.45), y - Inches(0.02), Inches(4), Inches(0.35),
                 title, font_size=17, color=COLOR_DARK, bold=True)
    # Mô tả
    add_text_box(slide, Inches(5.5), y + Inches(0.02), Inches(7), Inches(0.35),
                 desc, font_size=14, color=COLOR_GRAY)


# ============================================================
# SLIDE 3: TỔNG QUAN ĐỀ TÀI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "01. TỔNG QUAN ĐỀ TÀI")

# Mục tiêu
add_section_title(slide, Inches(0.6), Inches(1.4), "🎯 Mục tiêu đề tài")

tf = add_info_card(slide, Inches(0.6), Inches(1.95), Inches(5.8), Inches(2.6),
                    "Yêu cầu bài toán", "", COLOR_PRIMARY)
bullets = [
    "• Xây dựng phần mềm quản lý sinh viên hoàn chỉnh",
    "• Sử dụng cấu trúc dữ liệu Mảng động (Dynamic Array)",
    "• Giao diện đồ họa web thân thiện người dùng",
    "• Thực hiện đầy đủ CRUD: Thêm, Đọc, Sửa, Xóa",
    "• Phân tích thuật toán và đánh giá độ phức tạp",
]
add_bullet_slide_content(tf, bullets, font_size=14, spacing=Pt(4))

# Công nghệ
tf2 = add_info_card(slide, Inches(6.8), Inches(1.95), Inches(5.8), Inches(2.6),
                     "Công nghệ sử dụng", "", COLOR_TEAL)
bullets2 = [
    "• Ngôn ngữ: JavaScript (ES6+)",
    "• Giao diện: HTML5 + CSS3 (Responsive)",
    "• Kiến trúc: MVC (Model-View-Controller)",
    "• Lưu trữ: LocalStorage (JSON Serialization)",
    "• Xuất dữ liệu: SheetJS (XLSX)",
]
add_bullet_slide_content(tf2, bullets2, font_size=14, spacing=Pt(4))

# Phạm vi chức năng
add_section_title(slide, Inches(0.6), Inches(4.8), "📌 Phạm vi chức năng")

funcs = [
    ("Thêm mới", "Thêm sinh viên với kiểm tra\ntrùng mã SV", COLOR_SUCCESS),
    ("Tìm kiếm", "Tìm kiếm tuyến tính theo\nmã SV, họ tên, lớp", COLOR_SECONDARY),
    ("Sắp xếp", "TimSort theo mã, tên,\nđiểm TB, lớp", COLOR_PURPLE),
    ("Lọc & Thống kê", "Lọc theo lớp, thống kê\nđiểm TB, đạt/không đạt", COLOR_TEAL),
    ("Cập nhật / Xóa", "Sửa thông tin, xóa với\nxác nhận", COLOR_DANGER),
    ("Xuất Excel", "Xuất danh sách ra file\n.xlsx", COLOR_ACCENT),
]

for i, (title, desc, color) in enumerate(funcs):
    x = Inches(0.6) + Inches(i * 2.08)
    card = add_shape_bg(slide, x, Inches(5.35), Inches(1.9), Inches(1.7), 
                        RGBColor(0xF8, 0xF9, 0xFA))
    add_rect(slide, x, Inches(5.35), Inches(1.9), Inches(0.05), color)
    add_text_box(slide, x + Inches(0.1), Inches(5.45), Inches(1.7), Inches(0.35),
                 title, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(5.8), Inches(1.7), Inches(1.0),
                 desc, font_size=11, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4: KIẾN TRÚC HỆ THỐNG
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "02. KIẾN TRÚC HỆ THỐNG - MÔ HÌNH MVC")

# MVC diagram
mvc_items = [
    ("MODEL", "Dữ liệu & Logic", "MangDong.js\nSinhVien.js\ndata.js", 
     COLOR_SUCCESS, Inches(1.0)),
    ("VIEW", "Giao diện hiển thị", "index.html\nstyles.css", 
     COLOR_SECONDARY, Inches(5.2)),
    ("CONTROLLER", "Điều khiển logic", "BoQuanLySinhVien.js\nDieuKhienGiaoDien.js\nmain.js", 
     COLOR_PURPLE, Inches(9.4)),
]

for title, subtitle, files, color, x in mvc_items:
    # Card
    card = add_shape_bg(slide, x, Inches(1.5), Inches(3.2), Inches(3.0), 
                        RGBColor(0xF8, 0xF9, 0xFA))
    # Top bar
    add_rect(slide, x, Inches(1.5), Inches(3.2), Inches(0.6), color)
    add_text_box(slide, x, Inches(1.55), Inches(3.2), Inches(0.5),
                 title, font_size=22, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Subtitle
    add_text_box(slide, x + Inches(0.2), Inches(2.25), Inches(2.8), Inches(0.4),
                 subtitle, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Files
    add_text_box(slide, x + Inches(0.2), Inches(2.7), Inches(2.8), Inches(1.5),
                 files, font_size=14, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)

# Arrows between cards (simple text arrows)
add_text_box(slide, Inches(4.2), Inches(2.6), Inches(1), Inches(0.5),
             "⟶", font_size=36, color=COLOR_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(8.4), Inches(2.6), Inches(1), Inches(0.5),
             "⟶", font_size=36, color=COLOR_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

# Pipeline illustration
add_section_title(slide, Inches(0.6), Inches(4.8), "🔄 Pipeline xử lý dữ liệu hiển thị")

pipeline_steps = ["Dữ liệu gốc", "Lọc theo lớp", "Tìm kiếm", "Sắp xếp", "Phân trang", "Hiển thị"]
pipeline_complexity = ["O(n)", "O(n)", "O(n)", "O(n log n)", "O(k)", "O(k)"]

for i, (step, comp) in enumerate(zip(pipeline_steps, pipeline_complexity)):
    x = Inches(0.6) + Inches(i * 2.1)
    box = add_shape_bg(slide, x, Inches(5.4), Inches(1.7), Inches(1.1), 
                       RGBColor(0xE8, 0xEA, 0xF6))
    add_text_box(slide, x, Inches(5.45), Inches(1.7), Inches(0.4),
                 step, font_size=13, color=COLOR_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(5.9), Inches(1.7), Inches(0.4),
                 comp, font_size=16, color=COLOR_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    if i < len(pipeline_steps) - 1:
        add_text_box(slide, x + Inches(1.7), Inches(5.65), Inches(0.4), Inches(0.4),
                     "→", font_size=22, color=COLOR_SECONDARY, bold=True, alignment=PP_ALIGN.CENTER)

# Tổng pipeline
add_text_box(slide, Inches(0.6), Inches(6.65), Inches(12), Inches(0.5),
             "Độ phức tạp tổng pipeline: O(n log n) — bước Sắp xếp chiếm ưu thế",
             font_size=15, color=COLOR_DANGER, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5: CẤU TRÚC DỮ LIỆU MẢNG ĐỘNG
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "03. CẤU TRÚC DỮ LIỆU - MẢNG ĐỘNG (Dynamic Array)")

# Khái niệm
add_section_title(slide, Inches(0.6), Inches(1.4), "📖 Khái niệm & Nguyên lý")

tf = add_info_card(slide, Inches(0.6), Inches(1.95), Inches(6), Inches(2.5),
                    "Định nghĩa", "", COLOR_PRIMARY)
bullets = [
    "• Mảng động: cấu trúc dữ liệu tuyến tính tự thay đổi kích thước",
    "• Khác mảng tĩnh: không cần khai báo kích thước cố định",
    "• Tự động mở rộng khi đầy, thu nhỏ khi sử dụng quá ít",
    "• Truy cập ngẫu nhiên O(1) nhờ bộ nhớ liên tục",
    "• Thuộc tính: duLieu[], kichThuoc, dungLuong, soLanThayDoi",
]
add_bullet_slide_content(tf, bullets, font_size=14, spacing=Pt(4))

# Doubling Strategy
tf2 = add_info_card(slide, Inches(6.9), Inches(1.95), Inches(5.8), Inches(2.5),
                     "Chiến lược Doubling", "", COLOR_ACCENT)
bullets2 = [
    "• Mở rộng: kichThuoc >= dungLuong → dungLuong × 2",
    "• Thu nhỏ: kichThuoc < dungLuong/4 → dungLuong ÷ 2",
    "• Điều kiện tối thiểu: dungLuong >= 4",
    "• Hệ số tải (Load Factor): kichThuoc / dungLuong",
    "• Mở rộng khi LF ≥ 1.0, thu nhỏ khi LF < 0.25",
]
add_bullet_slide_content(tf2, bullets2, font_size=14, spacing=Pt(4))

# Minh họa mở rộng
add_section_title(slide, Inches(0.6), Inches(4.7), "📊 Minh họa quá trình mở rộng")

demo_data = [
    ("Thêm SV 1-4:", "dungLuong = 4", "[■][■][■][■]", "Đầy 100%"),
    ("Thêm SV 5:", "dungLuong = 8", "[■][■][■][■][■][ ][ ][ ]", "Mở rộng! → 62.5%"),
    ("Thêm SV 6-8:", "dungLuong = 8", "[■][■][■][■][■][■][■][■]", "Đầy 100%"),
    ("Thêm SV 9:", "dungLuong = 16", "[■][■]...[■][ ][ ][ ][ ][ ][ ][ ]", "Mở rộng! → 56.25%"),
]

for i, (action, cap, visual, note) in enumerate(demo_data):
    y = Inches(5.25) + Inches(i * 0.5)
    color = COLOR_DANGER if "Mở rộng" in note else COLOR_DARK
    add_text_box(slide, Inches(0.8), y, Inches(2.2), Inches(0.4),
                 action, font_size=13, color=COLOR_PRIMARY, bold=True)
    add_text_box(slide, Inches(3.0), y, Inches(1.8), Inches(0.4),
                 cap, font_size=13, color=COLOR_GRAY)
    add_text_box(slide, Inches(5.0), y, Inches(4.5), Inches(0.4),
                 visual, font_size=13, color=COLOR_DARK, font_name='Consolas')
    add_text_box(slide, Inches(10.0), y, Inches(2.5), Inches(0.4),
                 note, font_size=13, color=color, bold=True)


# ============================================================
# SLIDE 6: PHÂN TÍCH KHẤU HAO (Amortized Analysis)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "03b. PHÂN TÍCH CHI PHÍ KHẤU HAO (Amortized Analysis)")

# Giải thích
add_section_title(slide, Inches(0.6), Inches(1.4), "📐 Chứng minh bằng phương pháp Tổng (Aggregate Method)")

tf = add_info_card(slide, Inches(0.6), Inches(2.0), Inches(12.2), Inches(1.6),
                    "Bài toán", "", COLOR_PRIMARY)
bullets = [
    "• Mỗi lần thêm phần tử: chi phí O(1) nếu không cần mở rộng",
    "• Khi cần mở rộng: phải sao chép toàn bộ n phần tử → O(n)",
    "• Câu hỏi: Chi phí trung bình cho mỗi lần thêm là bao nhiêu?",
]
add_bullet_slide_content(tf, bullets, font_size=15, spacing=Pt(4))

# Chứng minh
add_section_title(slide, Inches(0.6), Inches(3.9), "🔢 Tính toán")

code = """Tổng chi phí cho n lần thêm liên tiếp:
  = n (chi phí thêm thông thường)
  + (1 + 2 + 4 + 8 + ... + n)  (chi phí sao chép khi mở rộng)
  = n + (2n - 1)
  = 3n - 1

Chi phí trung bình mỗi lần thêm = (3n - 1) / n ≈ 3 = O(1)

→ Kết luận: Chi phí khấu hao (Amortized Cost) = O(1)"""

add_code_block(slide, Inches(0.6), Inches(4.4), Inches(7), Inches(2.7), code, font_size=15)

# Bảng so sánh
tf3 = add_info_card(slide, Inches(8.0), Inches(3.9), Inches(4.8), Inches(3.2),
                     "So sánh với mảng tĩnh", "", COLOR_DANGER)
bullets3 = [
    "Mảng tĩnh:",
    "  • Kích thước cố định, không thể mở rộng",
    "  • Lãng phí bộ nhớ nếu khai báo quá lớn",
    "  • Tràn mảng nếu khai báo quá nhỏ",
    "",
    "Mảng động:",
    "  • Tự động điều chỉnh kích thước",
    "  • Sử dụng bộ nhớ hiệu quả",
    "  • Chi phí thêm trung bình O(1)",
]
add_bullet_slide_content(tf3, bullets3, font_size=13, spacing=Pt(2))


# ============================================================
# SLIDE 7: THUẬT TOÁN THÊM PHẦN TỬ
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "04a. THUẬT TOÁN THÊM PHẦN TỬ VÀO CUỐI")

# Pseudocode
add_section_title(slide, Inches(0.6), Inches(1.4), "📝 Mã giả (Pseudocode)")

code = """FUNCTION them(phanTu):
    IF kichThuoc >= dungLuong THEN
        _thayDoiKichThuoc(dungLuong * 2)  // Mở rộng gấp đôi
    END IF
    duLieu[kichThuoc] = phanTu            // Gán phần tử
    kichThuoc = kichThuoc + 1             // Tăng kích thước
    RETURN kichThuoc
END FUNCTION"""

add_code_block(slide, Inches(0.6), Inches(1.95), Inches(7), Inches(2.4), code, font_size=14)

# Thêm sinh viên (logic)
add_section_title(slide, Inches(8.0), Inches(1.4), "Thêm Sinh viên")

code2 = """FUNCTION themSinhVien(sinhVien):
    // B1: Kiểm tra trùng mã → O(n)
    viTriTrung = timViTri(sv.maSV == sinhVien.maSV)
    IF viTriTrung != -1 THEN
        THROW "Mã SV đã tồn tại!"
    END IF
    // B2: Thêm vào mảng → O(1)*
    mangDong.them(sinhVien)
    // B3: Lưu localStorage
    luuVaoBoNhoCucBo()
END FUNCTION"""

add_code_block(slide, Inches(8.0), Inches(1.95), Inches(4.8), Inches(2.65), code2, font_size=12)

# Bảng độ phức tạp
add_section_title(slide, Inches(0.6), Inches(4.6), "⏱️ Độ phức tạp")

data = [
    ["Thao tác", "Tốt nhất", "Xấu nhất", "Trung bình (Amortized)"],
    ["them() - không mở rộng", "O(1)", "O(1)", "O(1)"],
    ["them() - có mở rộng", "O(n)", "O(n)", "—"],
    ["them() - tổng thể", "O(1)", "O(n)", "O(1) ← khấu hao"],
    ["themSinhVien()", "O(n)", "O(n)", "O(n) ← do kiểm tra trùng"],
]

add_complexity_table(slide, Inches(0.6), Inches(5.15), data,
                     [Inches(3.5), Inches(2.2), Inches(2.2), Inches(3.5)])


# ============================================================
# SLIDE 8: THUẬT TOÁN XÓA PHẦN TỬ
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "04b. THUẬT TOÁN XÓA PHẦN TỬ TẠI VỊ TRÍ")

add_section_title(slide, Inches(0.6), Inches(1.4), "📝 Mã giả (Pseudocode)")

code = """FUNCTION xoaTaiViTri(viTri):
    IF viTri < 0 OR viTri >= kichThuoc THEN
        THROW "Vị trí ngoài phạm vi!"
    END IF
    
    phanTuDaXoa = duLieu[viTri]
    
    // Dịch chuyển phần tử sang trái → O(n)
    FOR i = viTri TO kichThuoc - 2:
        duLieu[i] = duLieu[i + 1]
    END FOR
    
    kichThuoc = kichThuoc - 1
    
    // Thu nhỏ nếu dùng < 25% dung lượng
    IF kichThuoc < dungLuong / 4 AND dungLuong > 4 THEN
        _thayDoiKichThuoc(dungLuong / 2)
    END IF
    
    RETURN phanTuDaXoa
END FUNCTION"""

add_code_block(slide, Inches(0.6), Inches(1.95), Inches(7.5), Inches(4.25), code, font_size=13)

# Minh họa dịch chuyển
tf2 = add_info_card(slide, Inches(8.5), Inches(1.95), Inches(4.3), Inches(4.25),
                     "Minh họa xóa tại vị trí 2", "", COLOR_DANGER)
bullets = [
    "Trước: [A][B][C][D][E]",
    "",
    "B1: Xóa phần tử C tại vị trí 2",
    "B2: Dịch D sang vị trí 2",
    "B3: Dịch E sang vị trí 3",
    "",
    "Sau:  [A][B][D][E][ ]",
    "",
    "Số lần dịch: n - viTri - 1",
    "  Xóa đầu → dịch n-1 lần",
    "  Xóa cuối → dịch 0 lần",
]
add_bullet_slide_content(tf2, bullets, font_size=13, spacing=Pt(2), font_name='Consolas')

# Bảng phức tạp
add_section_title(slide, Inches(0.6), Inches(6.4), "⏱️ Độ phức tạp")

data = [
    ["Vị trí xóa", "Tốt nhất", "Xấu nhất", "Trung bình"],
    ["Cuối mảng", "O(1)", "O(1)", "O(1)"],
    ["Đầu mảng", "O(n)", "O(n)", "O(n)"],
    ["Vị trí bất kỳ", "O(1)", "O(n)", "O(n)"],
]

add_complexity_table(slide, Inches(0.6), Inches(6.8), data,
                     [Inches(2.5), Inches(2), Inches(2), Inches(2)])


# ============================================================
# SLIDE 9: THUẬT TOÁN THAY ĐỔI KÍCH THƯỚC (RESIZE)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "04c. THUẬT TOÁN THAY ĐỔI KÍCH THƯỚC (Resize)")

add_section_title(slide, Inches(0.6), Inches(1.4), "📝 Mã giả (Pseudocode)")

code = """FUNCTION _thayDoiKichThuoc(dungLuongMoi):
    mangMoi = NEW Array(dungLuongMoi)
    
    // Sao chép dữ liệu → O(n)
    FOR i = 0 TO kichThuoc - 1:
        mangMoi[i] = duLieu[i]
    END FOR
    
    duLieu = mangMoi
    dungLuong = dungLuongMoi
    soLanThayDoi = soLanThayDoi + 1
END FUNCTION"""

add_code_block(slide, Inches(0.6), Inches(1.95), Inches(7), Inches(2.8), code, font_size=14)

# Chi tiết
tf2 = add_info_card(slide, Inches(8.0), Inches(1.95), Inches(4.8), Inches(2.8),
                     "Đặc điểm", "", COLOR_PURPLE)
bullets = [
    "• Cấp phát mảng mới với kích thước mới",
    "• Sao chép toàn bộ phần tử cũ sang mảng mới",
    "• Giải phóng mảng cũ (GC tự động)",
    "• Độ phức tạp: O(n) thời gian, O(n) không gian",
    "",
    "Khi nào gọi:",
    "  • Mở rộng: kichThuoc >= dungLuong",
    "  • Thu nhỏ: kichThuoc < dungLuong/4",
    "  • Tối thiểu: dungLuong >= 4",
]
add_bullet_slide_content(tf2, bullets, font_size=13, spacing=Pt(2))

# Bảng minh họa resize
add_section_title(slide, Inches(0.6), Inches(5.0), "📊 Lịch sử Resize khi thêm liên tục 20 phần tử")

data = [
    ["Lần resize", "Phần tử thứ", "Dung lượng cũ", "Dung lượng mới", "Số phần tử sao chép"],
    ["#1", "5", "4", "8", "4"],
    ["#2", "9", "8", "16", "8"],
    ["#3", "17", "16", "32", "16"],
]

add_complexity_table(slide, Inches(0.6), Inches(5.5), data,
                     [Inches(2), Inches(2), Inches(2.2), Inches(2.2), Inches(3)])

add_text_box(slide, Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
             "→ Chỉ 3 lần resize cho 20 lần thêm! Tổng chi phí sao chép = 4 + 8 + 16 = 28, trung bình = 28/20 ≈ 1.4 = O(1)",
             font_size=14, color=COLOR_DANGER, bold=True)


# ============================================================
# SLIDE 10: TÌM KIẾM TUYẾN TÍNH
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "05. THUẬT TOÁN TÌM KIẾM TUYẾN TÍNH (Linear Search)")

add_section_title(slide, Inches(0.6), Inches(1.4), "📝 Mã giả - Tìm vị trí theo điều kiện")

code = """FUNCTION timViTri(dieuKien):
    FOR i = 0 TO kichThuoc - 1:
        IF dieuKien(duLieu[i]) == TRUE THEN
            RETURN i          // Tìm thấy!
        END IF
    END FOR
    RETURN -1                 // Không tìm thấy
END FUNCTION"""

add_code_block(slide, Inches(0.6), Inches(1.95), Inches(6.5), Inches(2.2), code, font_size=14)

# Tìm kiếm sinh viên
add_section_title(slide, Inches(7.5), Inches(1.4), "Tìm kiếm trên nhiều trường")

code2 = """FUNCTION timKiemSinhVien(tuKhoa):
    tuKhoa = toLowerCase(tuKhoa)
    
    RETURN loc(sv =>
        sv.maSV    chứa tuKhoa  OR
        sv.hoTen   chứa tuKhoa  OR
        sv.tenLop  chứa tuKhoa
    )
END FUNCTION

// Sử dụng String.includes() → O(m)
// Tổng: O(n × m) ≈ O(n)"""

add_code_block(slide, Inches(7.5), Inches(1.95), Inches(5.3), Inches(2.8), code2, font_size=12)

# So sánh
add_section_title(slide, Inches(0.6), Inches(4.5), "⚖️ So sánh phương pháp tìm kiếm")

data = [
    ["Phương pháp", "Yêu cầu", "Độ phức tạp", "Phù hợp?"],
    ["Tìm kiếm tuyến tính", "Không yêu cầu sắp xếp", "O(n)", "✅ Đang sử dụng"],
    ["Tìm kiếm nhị phân", "Dữ liệu phải sắp xếp", "O(log n)", "❌ Không phù hợp"],
    ["Bảng băm (Hash Table)", "Cần bộ nhớ thêm", "O(1) trung bình", "❌ Phức tạp hơn"],
]

add_complexity_table(slide, Inches(0.6), Inches(5.05), data,
                     [Inches(3), Inches(3), Inches(2.5), Inches(3)])

tf3 = add_info_card(slide, Inches(0.6), Inches(6.3), Inches(12.2), Inches(0.85),
                     "Lý do chọn tìm kiếm tuyến tính", "", COLOR_TEAL)
bullets = [
    "• Dữ liệu không được sắp xếp sẵn → không dùng được nhị phân  •  Cần tìm trên nhiều trường (mã, tên, lớp) → tuyến tính linh hoạt hơn  •  Số lượng SV thực tế nhỏ (< 1000) → O(n) đủ nhanh"
]
add_bullet_slide_content(tf3, bullets, font_size=12, spacing=Pt(2))


# ============================================================
# SLIDE 11: THUẬT TOÁN SẮP XẾP (TimSort)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "06. THUẬT TOÁN SẮP XẾP - TimSort")

# Giới thiệu
add_section_title(slide, Inches(0.6), Inches(1.4), "📖 Giới thiệu TimSort")

tf = add_info_card(slide, Inches(0.6), Inches(1.95), Inches(6), Inches(2.5),
                    "Đặc điểm", "", COLOR_PURPLE)
bullets = [
    "• Thuật toán lai (Hybrid): Merge Sort + Insertion Sort",
    "• Phát minh: Tim Peters (2002)",
    "• Sử dụng trong: Python, Java, JavaScript Engine (V8)",
    "• Ổn định (Stable Sort): giữ thứ tự phần tử bằng nhau",
    "• Tận dụng các đoạn con đã sắp xếp (natural runs)",
    "• Thuật toán mặc định của Array.prototype.sort()",
]
add_bullet_slide_content(tf, bullets, font_size=14, spacing=Pt(4))

# Cách hoạt động
tf2 = add_info_card(slide, Inches(6.9), Inches(1.95), Inches(5.8), Inches(2.5),
                     "Nguyên lý hoạt động", "", COLOR_SECONDARY)
bullets2 = [
    "B1: Chia mảng thành các \"run\" (đoạn con đã sắp xếp)",
    "B2: Nếu run < minrun (thường 32-64):",
    "      → Dùng Insertion Sort mở rộng run",
    "B3: Merge các run lại theo Merge Sort",
    "B4: Tối ưu: galloping mode cho merge nhanh",
    "",
    "Ưu điểm: Rất nhanh với dữ liệu gần sắp xếp",
]
add_bullet_slide_content(tf2, bullets2, font_size=13, spacing=Pt(3))

# Implement trong hệ thống  
add_section_title(slide, Inches(0.6), Inches(4.7), "💻 Triển khai trong hệ thống")

code = """// Sắp xếp theo nhiều trường với localeCompare cho tiếng Việt
daSapXep.sort((a, b) => {
    switch (truongSapXep) {
        case 'hoTen': return a.hoTen.localeCompare(b.hoTen, 'vi'); // Tiếng Việt
        case 'diemTB': return a.diemTB - b.diemTB;                 // Số
        case 'maSV':   return a.maSV.localeCompare(b.maSV);        // Chuỗi
    }
    return sapXepTangDan ? ketQua : -ketQua;  // Tăng/Giảm dần
});"""

add_code_block(slide, Inches(0.6), Inches(5.25), Inches(8), Inches(2.0), code, font_size=13)

# Bảng phức tạp
tf3 = add_info_card(slide, Inches(8.9), Inches(4.7), Inches(3.9), Inches(2.55),
                     "Độ phức tạp TimSort", "", COLOR_DANGER)
bullets3 = [
    "Tốt nhất:    O(n)",
    "Trung bình:  O(n log n)",
    "Xấu nhất:    O(n log n)",
    "Không gian:   O(n)",
    "",
    "localeCompare('vi'):",
    "→ Sắp xếp đúng: Ă, Â, Đ, Ê, Ô, Ơ, Ư",
]
add_bullet_slide_content(tf3, bullets3, font_size=13, spacing=Pt(2), font_name='Consolas')


# ============================================================
# SLIDE 12: THUẬT TOÁN LỌC & THỐNG KÊ
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "07. THUẬT TOÁN LỌC DỮ LIỆU & THỐNG KÊ")

# Lọc dữ liệu
add_section_title(slide, Inches(0.6), Inches(1.4), "🔍 Thuật toán Lọc (Filter)")

code = """FUNCTION loc(dieuKien):
    ketQua = []                      // Mảng kết quả rỗng
    FOR i = 0 TO kichThuoc - 1:
        IF dieuKien(duLieu[i]) THEN
            ketQua.push(duLieu[i])   // Thêm phần tử thỏa mãn
        END IF
    END FOR
    RETURN ketQua                    // O(n) thời gian, O(k) không gian
END FUNCTION"""

add_code_block(slide, Inches(0.6), Inches(1.95), Inches(7.5), Inches(2.2), code, font_size=13)

# Các hàm lọc cụ thể
tf2 = add_info_card(slide, Inches(8.5), Inches(1.4), Inches(4.3), Inches(2.75),
                     "Ứng dụng trong hệ thống", "", COLOR_TEAL)
bullets = [
    "• locTheoLopHoc(tenLop): lọc SV theo lớp",
    "• demSinhVienDat(): đếm SV có điểm ≥ 5.0",
    "• demSinhVienChuaDat(): đếm SV có điểm < 5.0",
    "• timKiemSinhVien(): lọc theo từ khóa",
    "",
    "Tất cả đều dùng thuật toán lọc O(n)",
]
add_bullet_slide_content(tf2, bullets, font_size=13, spacing=Pt(3))

# Thống kê
add_section_title(slide, Inches(0.6), Inches(4.4), "📊 Thuật toán Thống kê")

code2 = """// Tính điểm trung bình chung → O(n)
FUNCTION tinhDiemTrungBinhChung():
    tong = 0
    FOR i = 0 TO kichThuoc - 1:
        tong += layGiaTri(i).diemTB
    END FOR
    RETURN tong / kichThuoc
END FUNCTION

// Danh sách lớp duy nhất → Set + Sort
FUNCTION layDanhSachLopDuyNhat():
    danhSachLop = NEW Set()          // Bảng băm → O(1) thêm
    FOR i = 0 TO kichThuoc - 1:
        danhSachLop.add(tenLop)      // Tự loại trùng
    END FOR
    RETURN Array.from(danhSachLop).sort()  // O(k log k)
END FUNCTION"""

add_code_block(slide, Inches(0.6), Inches(4.95), Inches(8), Inches(2.35), code2, font_size=12)

# Set explanation
tf3 = add_info_card(slide, Inches(8.9), Inches(4.4), Inches(3.9), Inches(2.9),
                     "Cấu trúc Set", "", COLOR_ACCENT)
bullets = [
    "• Set: tập hợp không trùng lặp",
    "• Nội bộ sử dụng bảng băm (Hash Table)",
    "• Thêm phần tử: O(1) trung bình",
    "• Kiểm tra tồn tại: O(1) trung bình",
    "",
    "Độ phức tạp tổng:",
    "  Duyệt n SV: O(n)",
    "  Sort k lớp: O(k log k)",
    "  Tổng: O(n + k log k) ≈ O(n)",
]
add_bullet_slide_content(tf3, bullets, font_size=12, spacing=Pt(2))


# ============================================================
# SLIDE 13: PHÂN TRANG
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "08. THUẬT TOÁN PHÂN TRANG (Pagination)")

add_section_title(slide, Inches(0.6), Inches(1.4), "📝 Mã giả")

code = """FUNCTION layDuLieuTheoTrang(danhSach):
    batDau = (trangHienTai - 1) × soDongMoiTrang
    ketThuc = batDau + soDongMoiTrang
    RETURN danhSach[batDau .. ketThuc]   // Array.slice()
END FUNCTION

FUNCTION layTongSoTrang(tongSoPhanTu):
    RETURN CEIL(tongSoPhanTu / soDongMoiTrang)
END FUNCTION

// Ví dụ: 25 sinh viên, 10 dòng/trang
// Trang 1: SV 0-9, Trang 2: SV 10-19, Trang 3: SV 20-24
// Tổng: CEIL(25/10) = 3 trang"""

add_code_block(slide, Inches(0.6), Inches(1.95), Inches(7.5), Inches(3.2), code, font_size=14)

# Đặc điểm
tf2 = add_info_card(slide, Inches(8.5), Inches(1.95), Inches(4.3), Inches(3.2),
                     "Đặc điểm", "", COLOR_SECONDARY)
bullets = [
    "• Dùng Array.slice() — không thay đổi mảng gốc",
    "• Chỉ tạo bản sao k phần tử (k = soDongMoiTrang)",
    "",
    "Độ phức tạp:",
    "  • Thời gian: O(k) với k = số dòng/trang",
    "  • Không gian: O(k)",
    "",
    "Cấu hình mặc định:",
    "  • soDongMoiTrang: 10",
    "  • Tùy chọn: 5, 10, 20 dòng/trang",
    "  • Hỗ trợ nút Previous/Next",
]
add_bullet_slide_content(tf2, bullets, font_size=13, spacing=Pt(2))

# Xếp loại
add_section_title(slide, Inches(0.6), Inches(5.4), "📋 Thuật toán Xếp loại học lực — O(1)")

code3 = """FUNCTION layXepLoai():
    IF diemTB >= 9.0 THEN RETURN "Xuất sắc"
    IF diemTB >= 8.0 THEN RETURN "Giỏi"
    IF diemTB >= 6.5 THEN RETURN "Khá"
    IF diemTB >= 5.0 THEN RETURN "Trung bình"
    IF diemTB >= 3.5 THEN RETURN "Yếu"
    RETURN "Kém"   // Tối đa 6 phép so sánh → O(1)
END FUNCTION"""

add_code_block(slide, Inches(0.6), Inches(5.9), Inches(7.5), Inches(1.45), code3, font_size=12)

tf4 = add_info_card(slide, Inches(8.5), Inches(5.4), Inches(4.3), Inches(1.95),
                     "Lưu trữ dữ liệu", "", COLOR_SUCCESS)
bullets2 = [
    "• LocalStorage: JSON Serialization/Deserialization",
    "• Lưu: Object → JSON.stringify() → Storage",
    "• Tải: Storage → JSON.parse() → Object",
    "• Độ phức tạp: O(n) cho cả lưu và tải",
]
add_bullet_slide_content(tf4, bullets2, font_size=12, spacing=Pt(2))


# ============================================================
# SLIDE 14: BẢNG TỔNG HỢP ĐỘ PHỨC TẠP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "09. BẢNG TỔNG HỢP ĐỘ PHỨC TẠP THUẬT TOÁN")

data = [
    ["Thuật toán / Thao tác", "Tốt nhất", "Trung bình", "Xấu nhất", "Không gian"],
    ["Thêm phần tử (them)", "O(1)", "O(1)*", "O(n)", "O(n)"],
    ["Xóa phần tử (xoaTaiViTri)", "O(1)", "O(n)", "O(n)", "O(1)"],
    ["Truy cập (layGiaTri)", "O(1)", "O(1)", "O(1)", "O(1)"],
    ["Cập nhật (ganGiaTri)", "O(1)", "O(1)", "O(1)", "O(1)"],
    ["Tìm kiếm tuyến tính (timViTri)", "O(1)", "O(n)", "O(n)", "O(1)"],
    ["Tìm kiếm sinh viên", "O(n)", "O(n)", "O(n)", "O(k)"],
    ["Sắp xếp TimSort", "O(n)", "O(n log n)", "O(n log n)", "O(n)"],
    ["Lọc dữ liệu (loc)", "O(n)", "O(n)", "O(n)", "O(k)"],
    ["Thay đổi kích thước (resize)", "O(n)", "O(n)", "O(n)", "O(n)"],
    ["Thống kê (điểm TB)", "O(n)", "O(n)", "O(n)", "O(1)"],
    ["Danh sách lớp duy nhất", "O(n)", "O(n)", "O(n)", "O(k)"],
    ["Phân trang (slice)", "O(k)", "O(k)", "O(k)", "O(k)"],
    ["Xếp loại học lực", "O(1)", "O(1)", "O(1)", "O(1)"],
    ["Lưu/Tải localStorage", "O(n)", "O(n)", "O(n)", "O(n)"],
]

add_complexity_table(slide, Inches(0.4), Inches(1.4), data,
                     [Inches(4), Inches(2), Inches(2.2), Inches(2.2), Inches(2)])

add_text_box(slide, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
             "* Chi phí khấu hao (Amortized) — k là số kết quả trả về, n là tổng số phần tử",
             font_size=13, color=COLOR_GRAY, bold=False)


# ============================================================
# SLIDE 15: KẾT LUẬN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "10. KẾT LUẬN & HƯỚNG PHÁT TRIỂN")

# Kết quả đạt được
add_section_title(slide, Inches(0.6), Inches(1.4), "✅ Kết quả đạt được")

tf = add_info_card(slide, Inches(0.6), Inches(1.95), Inches(6), Inches(3.2),
                    "Về lý thuyết", "", COLOR_SUCCESS)
bullets = [
    "• Hiểu rõ cấu trúc dữ liệu Mảng động",
    "• Nắm vững chiến lược Doubling Strategy",
    "• Phân tích chi phí khấu hao (Amortized Analysis)",
    "• Đánh giá đúng độ phức tạp các thuật toán",
    "• So sánh ưu/nhược với mảng tĩnh & DSLK",
    "• Hiểu nguyên lý TimSort trong JavaScript",
]
add_bullet_slide_content(tf, bullets, font_size=14, spacing=Pt(4))

tf2 = add_info_card(slide, Inches(6.9), Inches(1.95), Inches(5.8), Inches(3.2),
                     "Về thực hành", "", COLOR_SECONDARY)
bullets2 = [
    "• Xây dựng hoàn chỉnh hệ thống CRUD",
    "• Triển khai Mảng động từ đầu (class MangDong)",
    "• Giao diện đồ họa web responsive, hiện đại",
    "• Kiến trúc MVC rõ ràng, dễ bảo trì",
    "• Hỗ trợ: tìm kiếm, sắp xếp, lọc, phân trang",
    "• Xuất dữ liệu Excel, lưu trữ LocalStorage",
]
add_bullet_slide_content(tf2, bullets2, font_size=14, spacing=Pt(4))

# Hướng phát triển
add_section_title(slide, Inches(0.6), Inches(5.4), "🚀 Hướng phát triển")

improvements = [
    ("Tìm kiếm nhị phân", "Thêm index sắp xếp\ncho tìm kiếm O(log n)", COLOR_PRIMARY),
    ("Backend API", "Kết nối CSDL thực\n(MySQL, MongoDB)", COLOR_SECONDARY),
    ("Bảng băm", "Hash Table cho tra cứu\nmã SV nhanh O(1)", COLOR_PURPLE),
    ("Đồ thị thống kê", "Chart.js để trực quan\nhóa dữ liệu", COLOR_TEAL),
    ("Xác thực", "Đăng nhập, phân quyền\nngười dùng", COLOR_DANGER),
]

for i, (title, desc, color) in enumerate(improvements):
    x = Inches(0.6) + Inches(i * 2.5)
    card = add_shape_bg(slide, x, Inches(5.95), Inches(2.3), Inches(1.35),
                        RGBColor(0xF8, 0xF9, 0xFA))
    add_rect(slide, x, Inches(5.95), Inches(2.3), Inches(0.05), color)
    add_text_box(slide, x + Inches(0.1), Inches(6.05), Inches(2.1), Inches(0.3),
                 title, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(6.35), Inches(2.1), Inches(0.85),
                 desc, font_size=11, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 16: CẢM ƠN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_PRIMARY)

add_rect(slide, Inches(0), Inches(0), Inches(0.15), prs.slide_height, COLOR_ACCENT)
add_shape_bg(slide, Inches(8.5), Inches(-1), Inches(6), Inches(3),
             RGBColor(0x28, 0x3C, 0x8E))
add_shape_bg(slide, Inches(9), Inches(5.5), Inches(5.5), Inches(3),
             RGBColor(0x28, 0x3C, 0x8E))

add_text_box(slide, Inches(0), Inches(2.0), prs.slide_width, Inches(1.5),
             "CẢM ƠN THẦY VÀ CÁC BẠN\nĐÃ LẮNG NGHE!",
             font_size=44, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(5.5), Inches(3.55), Inches(2.5), Inches(0.06), COLOR_ACCENT)

add_text_box(slide, Inches(0), Inches(4.0), prs.slide_width, Inches(0.8),
             "GVHD: TS. Nghiêm Văn Tính\nSVTH: Đặng Đình Đạt",
             font_size=22, color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(5.2), prs.slide_width, Inches(0.5),
             "Môn: Kỹ thuật Lập trình Nâng cao  |  Năm học 2025 - 2026",
             font_size=16, color=RGBColor(0x90, 0xCA, 0xF9), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(6.0), prs.slide_width, Inches(0.5),
             "❓ Câu hỏi & Thảo luận",
             font_size=20, color=COLOR_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# LƯU FILE
# ============================================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 
                           "ThuatToan_QuanLySinhVien.pptx")
prs.save(output_path)
print(f"✅ Đã tạo file thuyết trình thành công!")
print(f"📁 Đường dẫn: {output_path}")
print(f"📊 Tổng số slide: {len(prs.slides)}")
